import os
from dotenv import load_dotenv
import json
from openai import OpenAI, AzureOpenAI
from groq import Groq
from pptx import Presentation
from pptx.util import Pt

load_dotenv()

client = Groq(
    api_key=os.environ.get("GROQ_API_KEY"),
)

# print(chat_completion.choices[0].message.content)

# LLM
# client = AzureOpenAI(
#     azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
#     api_key=os.getenv("AZURE_OPENAI_API_KEY"),
#     api_version="2024-06-01",
# )

# Uncomment for LM Studio
#client = OpenAI(base_url="http://localhost:1234/v1", api_key="lm-studio")

# Prompt setup
query_json = """{
    "input_text": "[[content]]",
    "output_format": "json",
    "json_structure": {
    "slides": "{{presentation_slides}}"
    }
}"""

user_prompt = (
    "Generate a 10 slide presentation for the user guide: "
    + """User Guide: Creating a New Microsoft Outlook Account
    1. Access the Outlook sign-in page by using your preferred email client or web browser and visiting www.outlook.com.
    2. Click on the "Create Account" option located at the bottom of the sign-in form.
    3. Fill out the account creation form with your desired email address, password (ensure it's unique), personal information, security question, and accept the Microsoft Terms & Privacy Agreement by checking the corresponding box.
    4. Click on "Create Account" to initiate the process of creating a new Microsoft Outlook account.
    5. Check your email inbox for the confirmation email containing an activation link; click it upon receipt.
    6. After activating your account, log in and customize settings using www.outlook.com or through compatible email clients with POP3/IMAP access to complete setup tasks such as adding contacts, configuring preferences, etc.
    7. Comply with any additional security measures like two-factor authentication if required for enhanced protection of your account."""
    + ". Produce 50 to 60 words per slide. Each slide should have a {{header}}, {{content}}. "
    "The final slide should be a list of discussion questions. Return only the JSON data as specified below."
)

system_prompt = "You are a helpful assistant. Your task is to generate only the JSON data required by the user and nothing else. Do not include any additional text or explanations."

prompt = query_json.replace("[[content]]", user_prompt)

# Chat completions API call
completion = client.chat.completions.create(
    model="llama3-8b-8192",
    messages=[
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": prompt}
    ]
)

# Retrieve response
response_content = completion.choices[0].message.content
print(response_content)

# Parse the response as JSON
json_rsp = json.loads(response_content)


slide_data = json_rsp.get("slides")

prs = Presentation()

# Create slide.py
for slide in slide_data:
    slide_layout = prs.slide_layouts[1]
    new_slide = prs.slides.add_slide(slide_layout)

    if slide['header']:
        title = new_slide.shapes.title
        title.text = slide['header']

    if slide['content']:
        shapes = new_slide.shapes
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame

        content = slide['content']
        if isinstance(content, list):
            content = " ".join(content)

        p = tf.add_paragraph()
        run = p.add_run()
        run.text = content
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.name = "Calibri"

prs.save("output.pptx")
print("Presentation saved as output.pptx")
