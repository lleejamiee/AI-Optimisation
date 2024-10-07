import json
from dotenv import load_dotenv
import os
from pathlib import Path
from llama_index.core import Settings
from llama_index.llms.groq import Groq
from llama_parse import LlamaParse
from autogen import ConversableAgent
from pptx import Presentation
from llama_index.readers.web import SimpleWebPageReader

load_dotenv()

review_sys_prompt = """You are tasked with identifying and retaining only the sections of the document that need to 
be replaced based on the provided reference material. Your goal is to extract the content that is directly relevant 
for replacement without adding or modifying the existing text.

Instructions: 1. Review both the provided document chunk and reference material to determine which sections are 
directly relevant for replacement. 2. Extract the content that will be replaced, including any references or context 
that should be retained as part of the replacement. 3. Ensure that the extracted section remains as close as possible 
in length and context to the original content that will be replaced. 4. Do not include any additional content or 
modifications outside of what is necessary for the update. 5. Do not include any explanations, introductions, 
or additional text."""

generate_sys_prompt = """
You are tasked to modify the text that is given by the user based on the reference material.

Instructions:
1. Review all the materials that are provided to generate an updated text or a new content.
2. If there is nothing to update, display 'n/a'.
3. If the content needs to be updated, generate an updated text.
4. If the content should be added with new material, generate a new text.
5. If the content should be deleted, display 'remove'
6. Do not include any additional content.
7. Do not include any explanations, introductions, or additional text.
"""

create_sys_prompt = """You are tasked with creating a new PowerPoint presentation for the user based on the provided reference materials and the prompt input by the user.

**Instructions:**
1. Carefully review all reference materials provided.
2. Analyze the user's prompt to understand their requirements and expectations.
3. Create a PowerPoint presentation structured in the following JSON format:

```json
{
  "slides": [
    {
      "heading": "{heading}",
      "content": "{content}"
    }
  ]
}
4. Do not include any additional content.
5. Do not include any explanations, introductions, or additional text.
"""
class SlideGenerator:
    Settings.llm = Groq(model="llama-3.1-70b-versatile", api_key=os.getenv("GROQ_API_KEY"), temperature=0.4)

    config_list = [
        {
            "model": "llama-3.1-70b-versatile",
            "api_key": os.getenv("GROQ_API_KEY"),
            "api_type": "groq",
            "frequency_penalty": 0.5,
            "max_tokens": 8000,
            "presence_penalty": 0.2,
            "seed": 42,
            "temperature": 0,
            "top_p": 0.2
        }
    ]

    agent = ConversableAgent(
        "chatbot",
        llm_config={"config_list": config_list},
        code_execution_config=False,  # Turn off code execution, by default it is off.
        function_map=None,  # No registered functions, by default it is None.
        human_input_mode="NEVER",  # Never ask for human input.
    )

    def save_uploaded_file(file):
        save_folder = Path('uploaded_files/')
        save_folder.mkdir(parents=True, exist_ok=True)
        save_path = save_folder / file.name
        with open(save_path, mode='wb') as w:
            w.write(file.getvalue())

        return save_path

    def load_reference(file_path):
        reference = LlamaParse().load_data(file_path)

        return reference

    def load_webpage(url):
        document = SimpleWebPageReader(html_to_text=True).load_data([url])

        return document

    def generate_updated(outdated_file_path, reference):
        ppt = Presentation(outdated_file_path)

        contents_map = {}

        review_system_message = {"content": review_sys_prompt, "role": "system"}
        generate_system_message = {"content": generate_sys_prompt, "role": "system"}

        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue

                if shape.has_text_frame:
                    review_context = f""" Analyse the given content and reference material to determine if the 
                    content needs to be updated or new material should be added to the content. This is the content: 
                    {shape.text}
                    Here is the reference material to use for determining relevance:
                    {reference} If nothing needs to be changed, display 'no'. If the content needs to be updated or 
                    new material should be added, display 'yes'"""

                    review_user_message = {"content": review_context, "role": "user"}
                    reply = SlideGenerator.agent.generate_reply(messages=[review_system_message, review_user_message])

                    for paragraph in shape.text_frame.paragraphs:
                        if not paragraph.text.strip():
                            continue

                        contents_map[paragraph.text] = paragraph.text

                        if reply:
                            if reply['content'] == "yes":

                                generate_context = f""" Based on the given context and the reference material, update 
                                the provided sentence or generate a new content.
        
                                This is the context:
                                {shape.text}
                                Here is the reference material:
                                {reference} Here is the sentence that needs to be updated with new information from 
                                the reference material: {paragraph.text} If nothing needs to be updated or added, 
                                or if the sentence hasn't been provided, display 'n/a'. If the given sentence needs 
                                to be removed, display 'remove'."""

                                generate_user_message = {"content": generate_context, "role": "user"}
                                reply = SlideGenerator.agent.generate_reply(messages=[generate_system_message, generate_user_message])

                                if reply:
                                    if reply['content'] != "n/a" and reply['content'] != "remove":
                                        contents_map[paragraph.text] = reply['content']
                                        paragraph.text = paragraph.text.replace(paragraph.text, reply['content'])
                                    if reply['content'] == "remove":
                                        paragraph.text = paragraph.text.replace(paragraph.text, "")

        return ppt, contents_map

    def generate_new(reference_docs, prompt):
        generate_sys_message = {"content": create_sys_prompt, "role": "system"}

        context = f"""Based on the given user prompt and the reference materials, generate a new powerpoint. 
        This is the user prompt: {prompt}
        Here is the reference material: {reference_docs}"""
        generate_user_message = {"content": context, "role": "user"}

        reply = SlideGenerator.agent.generate_reply(messages=[generate_sys_message, generate_user_message])

        return reply

    def extract_json(reply):
        json_rsp = json.loads(reply['content'])
        slide_data = json_rsp.get("slides")

        return slide_data