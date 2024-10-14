from html import unescape
import re
from pptx import Presentation
from pptx.util import Inches, Pt
import streamlit as st


class SlideProcessor:
    # Method to compare differences between the before and after contents
    def compare_difference(**kwargs):
        before = []
        after = []

        before_list = list(st.session_state.contents_map.keys())
        after_list = list(st.session_state.contents_map.values())

        max_length = max(len(before_list), len(after_list))

        for i in range(max_length):
            before_text = before_list[i] if i < len(before_list) else ""
            after_text = after_list[i] if i < len(after_list) else ""

            if before_text == after_text:
                before.append(before_text)
                after.append(after_text)

            elif after_text == "":
                before.append(f'<span style="color: #FFCCCB;">{before_text}</span>')
                after.append(after_text)

            elif before_text == "":
                before.append(before_text)
                after.append(f'<span style="color: #90EE90;">{after_text}</span>')

            else:
                before.append(f'<span style="color: #FFCCCB;">{before_text}</span>')
                after.append(f'<span style="color: #90EE90;">{after_text}</span>')

        st.session_state.before = before
        st.session_state.after = after

    # Method to display the differences between the before and after contents
    def display_difference(before, after):
        col1, col2 = st.columns(2)

        with col1:
            st.subheader("Outdated User Guide")
            st.markdown("<br>".join(before), unsafe_allow_html=True)

        with col2:
            st.subheader("Updated User Guide")
            st.markdown("<br>".join(after), unsafe_allow_html=True)

    # Method to remove the HTML tags for editing
    def strip_html_tags(**kwargs):
        clean_text = ""
        for text in st.session_state.after:
            cleaned = re.sub(r"<[^>]+>", "", text)
            cleaned = unescape(cleaned)
            clean_text += cleaned + "\n"

        return clean_text

    # Method to allow users to edit the generated contents
    def edit_output(**kwargs):
        clean_text = SlideProcessor.strip_html_tags()

        edited_text = st.text_area(
            "Edit the content",
            height=800,
            value=clean_text,
            key="edit_area",
        )

        if st.button("Save"):
            SlideProcessor.save_edit(edited_text)
            st.rerun()

    # Method to save the edited text
    def save_edit(edited_text):
        st.session_state.edit_mode = False

        st.session_state.after = edited_text.split("\n")

        i = 1
        new_map = {}
        # Iterate through the before and after contents
        for before, after in zip(st.session_state.before, st.session_state.after):
            if not after.strip():  # Handle when the user added new lines
                new_map[f'New Line {i}'] = "\n"
                i += 1
            else:
                new_map[before] = after

        i = 1
        # Handle f the after content is longer than before content
        if len(st.session_state.after) > len(st.session_state.before):
            for extra_text in st.session_state.after[len(st.session_state.before):]:
                new_map[f'Extra {i}'] = extra_text
                i += 1

        # Handle if the before content is longer than the after content
        if len(st.session_state.before) > len(st.session_state.after):
            for extra_before_text in st.session_state.before[len(st.session_state.after):]:
                new_map[extra_before_text] = ""

        st.session_state.contents_map = new_map

        SlideProcessor.compare_difference()

    # Method to modify the outdated PowerPoint with new content
    def modify_powerpoint(**kwargs):
        extra_text = ""
        for key, value in st.session_state.contents_map.items():  # Iterate through the before and after
            for slide in st.session_state.ppt.slides:  # Iterate through the outdated slide
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        if key in paragraph.text:
                            if "New Line" in key:  # Handle if the user added new line
                                new_paragraph = shape.text_frame.add_paragraph()
                                new_paragraph.text = ""
                            if paragraph.text != value:  # Handle updated content
                                paragraph.text = value
            # When the length of the updated material is longer than outdated, store them separately
            if "Extra" in key:
                extra_text += value + "\n"

        # Append all the extra material into one page at the end of the PowerPoint document
        if extra_text:
            new_slide_layout = st.session_state.ppt.slide_layouts[6]
            new_slide = st.session_state.ppt.slides.add_slide(new_slide_layout)

            left = top = Inches(1)
            height = st.session_state.ppt.slide_height - 2 * top
            width = st.session_state.ppt.slide_width - 2 * left

            tx_box = new_slide.shapes.add_textbox(left, top, width, height)
            tf = tx_box.text_frame

            p = tf.add_paragraph()
            p.text = "New Content"
            p.font.size = Pt(40)
            p.font.bold = True

            p = tf.add_paragraph()
            p.text = extra_text
            p.font.size = Pt(12)

    # Method to convert JSON into a list
    def json_to_string(**kwargs):
        all_content = []

        for slide in st.session_state.slide_data:
            all_content.append(slide['heading'])
            all_content.append(slide['content'])

        return all_content

    # Method to create new PowerPoint document
    def create_slides(**kwargs):
        ppt = Presentation()

        for slide in st.session_state.slide_data:
            slide_layout = ppt.slide_layouts[6]
            new_slide = ppt.slides.add_slide(slide_layout)

            left = top = Inches(0.5)
            height = ppt.slide_height - 2 * top
            width = ppt.slide_width - 2 * left

            tx_box = new_slide.shapes.add_textbox(left, top, width, height)
            tf = tx_box.text_frame

            p = tf.add_paragraph()
            p.text = slide['heading']
            p.font.size = Pt(40)
            p.font.bold = True

            p = tf.add_paragraph()
            p.text = slide['content']
            p.font.size = Pt(12)

        return ppt
