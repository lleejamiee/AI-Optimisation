from dotenv import load_dotenv
import streamlit as st
import os
import re
from pptx.util import Inches, Pt
from pathlib import Path
from html import unescape
from llama_index.core import Settings
from llama_index.llms.groq import Groq
from llama_parse import LlamaParse
from autogen import ConversableAgent
from pptx import Presentation
from llama_index.readers.web import SimpleWebPageReader

load_dotenv()

review_sys_prompt = """You are tasked with identifying and retaining only the sections of the document that need to 
be replaced based on the provided reference material. Your goal is to extract the content that is directly relevant 
for replacement without adding or modifying the existing text.

Instructions: 1. Review both the provided document chunk and reference material to determine which sections are 
directly relevant for replacement. 2. Extract the content that will be replaced, including any references or context 
that should be retained as part of the replacement. 3. Ensure that the extracted section remains as close as possible 
in length and context to the original content that will be replaced. 4. Do not include any additional content or 
modifications outside of what is necessary for the update. 5. Do not include any explanations, introductions, 
or additional text."""

generate_sys_prompt = """
You are tasked to modify the text that is given by the user based on the reference material.

Instructions:
1. Review all the materials that are provided to generate an updated text or a new content.
2. If there is nothing to update, display 'n/a'.
3. If the content needs to be updated, generate an updated text.
4. If the content should be added with new material, generate a new text.
5. If the content should be deleted, display 'remove'
6. Do not include any additional content.
7. Do not include any explanations, introductions, or additional text.
"""
class SlideRag:
    Settings.llm = Groq(model="llama-3.1-70b-versatile", api_key=os.getenv("GROQ_API_KEY"), temperature=0.4)

    config_list = [
        {
            "model": "llama-3.1-70b-versatile",
            "api_key": os.getenv("GROQ_API_KEY"),
            "api_type": "groq",
            "frequency_penalty": 0.5,
            "max_tokens": 8000,
            "presence_penalty": 0.2,
            "seed": 42,
            "temperature": 0,
            "top_p": 0.2
        }
    ]

    agent = ConversableAgent(
        "chatbot",
        llm_config={"config_list": config_list},
        code_execution_config=False,  # Turn off code execution, by default it is off.
        function_map=None,  # No registered functions, by default it is None.
        human_input_mode="NEVER",  # Never ask for human input.
    )

    def save_uploaded_file(file):
        save_folder = Path('uploaded_files/')
        save_folder.mkdir(parents=True, exist_ok=True)
        save_path = save_folder / file.name
        with open(save_path, mode='wb') as w:
            w.write(file.getvalue())

        return save_path

    def load_reference(file_path):
        reference = LlamaParse().load_data(file_path)

        return reference

    def load_webpage(url):
        document = SimpleWebPageReader(html_to_text=True).load_data([url])

        return document

    def generate_updated(outdated_file_path, reference):
        ppt = Presentation(outdated_file_path)

        contents_map = {}

        review_system_message = {"content": review_sys_prompt, "role": "system"}
        generate_system_message = {"content": generate_sys_prompt, "role": "system"}

        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue

                if shape.has_text_frame:
                    review_context = f""" Analyse the given content and reference material to determine if the 
                    content needs to be updated or new material should be added to the content. This is the content: 
                    {shape.text}
                    Here is the reference material to use for determining relevance:
                    {reference} If nothing needs to be changed, display 'no'. If the content needs to be updated or 
                    new material should be added, display 'yes'"""

                    review_user_message = {"content": review_context, "role": "user"}
                    reply = SlideRag.agent.generate_reply(messages=[review_system_message, review_user_message])

                    for paragraph in shape.text_frame.paragraphs:
                        if not paragraph.text.strip():
                            continue

                        contents_map[paragraph.text] = paragraph.text

                        if reply:
                            if reply['content'] == "yes":

                                generate_context = f""" Based on the given context and the reference material, update 
                                the provided sentence or generate a new content.
        
                                This is the context:
                                {shape.text}
                                Here is the reference material:
                                {reference} Here is the sentence that needs to be updated with new information from 
                                the reference material: {paragraph.text} If nothing needs to be updated or added, 
                                or if the sentence hasn't been provided, display 'n/a'. If the given sentence needs 
                                to be removed, display 'remove'."""

                                generate_user_message = {"content": generate_context, "role": "user"}
                                reply = SlideRag.agent.generate_reply(messages=[generate_system_message, generate_user_message])

                                if reply:
                                    if reply['content'] != "n/a" and reply['content'] != "remove":
                                        contents_map[paragraph.text] = reply['content']
                                        paragraph.text = paragraph.text.replace(paragraph.text, reply['content'])
                                    if reply['content'] == "remove":
                                        paragraph.text = paragraph.text.replace(paragraph.text, "")

        return ppt, contents_map

    def compare_difference(**kwargs):
        before = []
        after = []

        before_list = list(st.session_state.contents_map.keys())
        after_list = list(st.session_state.contents_map.values())

        max_length = max(len(before_list), len(after_list))

        for i in range(max_length):
            before_text = before_list[i] if i < len(before_list) else ""
            after_text = after_list[i] if i < len(after_list) else ""

            if before_text == after_text:
                before.append(before_text)
                after.append(after_text)

            elif after_text == "":
                before.append(f'<span style="color: #FFCCCB;">{before_text}</span>')
                after.append(after_text)

            elif before_text == "":
                before.append(before_text)
                after.append(f'<span style="color: #90EE90;">{after_text}</span>')

            else:
                before.append(f'<span style="color: #FFCCCB;">{before_text}</span>')
                after.append(f'<span style="color: #90EE90;">{after_text}</span>')

        st.session_state.before = before
        st.session_state.after = after

    def display_difference(before, after):
        col1, col2 = st.columns(2)

        with col1:
            st.subheader("Outdated User Guide")
            st.markdown("<br>".join(before), unsafe_allow_html=True)

        with col2:
            st.subheader("Updated User Guide")
            st.markdown("<br>".join(after), unsafe_allow_html=True)

    def strip_html_tags(**kwargs):
        clean_text = ""
        for text in st.session_state.after:
            cleaned = re.sub(r"<[^>]+>", "", text)
            cleaned = unescape(cleaned)
            clean_text += cleaned + "\n"

        return clean_text


    def edit_output(**kwargs):
        clean_text = SlideRag.strip_html_tags()

        edited_text = st.text_area(
            "Edit the content",
            height=800,
            value=clean_text,
            key="edit_area",
        )

        if st.button("Save"):
            SlideRag.save_edit(edited_text)
            st.rerun()

    def save_edit(edited_text):
        st.session_state.edit_mode = False

        st.session_state.after = edited_text.split("\n")

        i = 1
        new_map = {}
        for before, after in zip(st.session_state.before, st.session_state.after):
            if not after.strip():
                new_map[f'New Line {i}'] = "\n"
                i += 1
            else:
                new_map[before] = after

        i = 1
        if len(st.session_state.after) > len(st.session_state.before):
            for extra_text in st.session_state.after[len(st.session_state.before):]:
                new_map[f'Extra {i}'] = extra_text
                i += 1

        if len(st.session_state.before) > len(st.session_state.after):
            for extra_before_text in st.session_state.before[len(st.session_state.after):]:
                new_map[extra_before_text] = ""

        st.session_state.contents_map = new_map

        SlideRag.compare_difference()

    def modify_powerpoint(**kwargs):
        extra_text = ""
        for key, value in st.session_state.contents_map.items():
            for slide in st.session_state.ppt.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        if key in paragraph.text:
                            if "New Line" in key:
                                new_paragraph = shape.text_frame.add_paragraph()
                                new_paragraph.text = ""
                            if paragraph.text != value:
                                paragraph.text = value
            if "Extra" in key:
                extra_text += value + "\n"

        if extra_text:
            new_slide_layout = st.session_state.ppt.slide_layouts[6]
            new_slide = st.session_state.ppt.slides.add_slide(new_slide_layout)

            left = top = Inches(1)
            height = st.session_state.ppt.slide_height - 2 * top
            width = st.session_state.ppt.slide_width - 2 * left

            tx_box = new_slide.shapes.add_textbox(left, top, width, height)
            tf = tx_box.text_frame

            p = tf.add_paragraph()
            p.text = "New Content"
            p.font.size = Pt(40)
            p.font.bold = True

            p = tf.add_paragraph()
            p.text = extra_text
            p.font.size = Pt(12)
