from urllib.request import urlopen
from pptx import Presentation
from bs4 import BeautifulSoup
from html import unescape
import streamlit as st
import pdfplumber
import requests
import docx2txt
import difflib
import certifi
import ssl
import re
import os
from dotenv import load_dotenv
from groq import Groq

load_dotenv()


class DocumentProcessor:
    client = Groq(
        api_key=os.environ.get("GROQ_API_KEY"),
    )

    def generate_text(prompt, outdated_guide, reference_material):
        history = [
            {"role": "system",
             "content": prompt},
            {"role": "user",
             "content": f"Outdated User Guide:\n{outdated_guide}\n\nReference Material:\n{reference_material}"}
        ]

        chat_completion = DocumentProcessor.client.chat.completions.create(
            messages=history,
            model="llama3-8b-8192",
        )

        generated_text = chat_completion.choices[0].message.content

        return generated_text

    def extract_text(file):
        if file.type == "text/plain":
            file_text = str(file.read(), "utf-8")
        elif file.type == "application/pdf":
            try:
                with pdfplumber.open(file) as pdf:
                    pages = pdf.pages[0]
                    file_text = pages.extract_text()
            except:
                return
        elif (
                file.type
                == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            file_text = docx2txt.process(file)
        elif (
                file.type
                == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            presentation = Presentation(file)
            file_text = ""

            for slide_number, slide in enumerate(presentation.slides):
                file_text += f"\nSlide {slide_number + 1}:\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        file_text += shape.text + "\n"

        return file_text

    def extract_url_text(url):
        try:
            resp = urlopen(
                url, context=ssl.create_default_context(cafile=certifi.where())
            )
            context = ssl._create_unverified_context()
            response = urlopen(url, context=context)
            html = response.read()

            soup = BeautifulSoup(html, "html.parser")
            title = soup.title.string

            headings = [
                h.get_text()
                for h in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p"])
            ]

            content = " ".join(headings)

            return content

        except requests.exceptions.RequestException as e:
            return f"Error: {e}"

    def handle_updated_guide(updated_guide):
        st.write(updated_guide)

        st.subheader("Download the Updated User Guide")

        # Create .docx file
        with open("updated_guide.docx", "w") as f:
            # Write from generated text
            f.write(updated_guide)

        # Allow app user to download .docx file(updated user guide file) containing generated text
        st.download_button("Download", updated_guide)

    def read_files(outdated_guide_file, reference_material_file):
        if outdated_guide_file is not None and reference_material_file is not None:
            outdated_guide_text = DocumentProcessor.extract_text(outdated_guide_file)
            reference_material_text = DocumentProcessor.extract_text(
                reference_material_file
            )

            return outdated_guide_text, reference_material_text
        else:
            st.error("Please upload both files.")
            return None, None

    def edit_output(updated_guide):
        # Display the text area for editing
        edited_text = st.text_area(
            "Edit the content",
            height=800,
            value=st.session_state.generated_text,
            key="edit_area",
        )

        # Save the changes
        if st.button("Save Changes"):
            st.session_state.generated_text = edited_text
            st.success("Changes saved successfully!")
            st.session_state.edit_mode = False
            st.rerun()

    def display_output(updated_guide):
        st.write(st.session_state.generated_text)

        if st.button("Edit this output"):
            st.session_state.edit_mode = True
            st.rerun()

    def display_output(outdated_guide, updated_guide):
        col1, col2 = st.columns(2)

        with col1:
            st.subheader("Outdated User Guide")
            st.markdown("<br>".join(outdated_guide), unsafe_allow_html=True)

        with col2:
            st.subheader("Updated User Guide")
            st.markdown("<br>".join(updated_guide), unsafe_allow_html=True)

        if st.button("Edit this output"):
            st.session_state.edit_mode = True
            st.rerun()

    def strip_html_tags(text):
        # Remove HTML tags
        clean_text = re.sub(r"<[^>]+>", "", text)
        # Unescape HTML entities
        clean_text = unescape(clean_text)
        return clean_text

    def compare_texts(text1, text2):
        text1_lines = text1.splitlines()
        text2_lines = text2.splitlines()

        differ = difflib.Differ()
        diff = list(differ.compare(text1_lines, text2_lines))

        text1_output = []
        text2_output = []

        for line in diff:
            if line.startswith("  "):  # Unchanged line
                text1_output.append(line[2:].rstrip())
                text2_output.append(line[2:].rstrip())
            elif line.startswith("- "):  # Line only in text1
                text1_output.append(
                    f'<span style="color: #FFCCCB;">{line[2:].rstrip()}</span>'
                )
            elif line.startswith("+ "):  # Line only in text2
                text2_output.append(
                    f'<span style="color: #90EE90;">{line[2:].rstrip()}</span>'
                )
            elif line.startswith("? "):  # Hint line, ignore
                continue

        return text1_output, text2_output

    def select_version(**kwargs):
        department, roles = st.columns(2)

        with department:
            option = st.selectbox(
                "Choose a department", ("ICT", "HR", "Marketing"), index=None
            )

        if option:
            with roles:
                if option == "ICT":
                    st.checkbox("Manager")
                    st.checkbox("IT Technician")
                    st.checkbox("Team Member")
                    st.checkbox("Graduate")
                    st.checkbox("Intern")
                elif option == "HR":
                    st.checkbox("Manager")
                    st.checkbox("Recruiter")
                    st.checkbox("Team Member")
                    st.checkbox("Graduate")
                    st.checkbox("Intern")
                elif option == "Marketing":
                    st.checkbox("Manager")
                    st.checkbox("Marketing Specialist")
                    st.checkbox("Content Creator")
                    st.checkbox("Graduate")
                    st.checkbox("Intern")
